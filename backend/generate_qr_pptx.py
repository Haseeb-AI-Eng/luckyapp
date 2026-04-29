#!/usr/bin/env python3
"""
Generate a PowerPoint presentation with a printable QR code for Lucky Draw registration
"""

import qrcode
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import io
from datetime import datetime

def generate_qr_pptx(registration_code="Z5G-12345"):
    """
    Generate a PowerPoint file with QR code
    
    Args:
        registration_code (str): The code to encode in QR (format: XXX-#####)
    """
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Add blank slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Set background to white
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Generate QR Code
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_H,
        box_size=10,
        border=2,
    )
    qr.add_data(registration_code)
    qr.make(fit=True)
    
    # Create QR image (white background, blue foreground)
    qr_img = qr.make_image(fill_color="rgb(30, 64, 175)", back_color="white")
    
    # Save QR to bytes
    qr_bytes = io.BytesIO()
    qr_img.save(qr_bytes, format='PNG')
    qr_bytes.seek(0)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.text = "Lucky Draw Registration"
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(54)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(30, 64, 175)
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.text = "Scan the QR Code Below to Register"
    subtitle_paragraph = subtitle_frame.paragraphs[0]
    subtitle_paragraph.font.size = Pt(24)
    subtitle_paragraph.font.color.rgb = RGBColor(100, 116, 139)
    subtitle_paragraph.alignment = PP_ALIGN.CENTER
    
    # Add QR Code image (centered)
    qr_width = Inches(4)
    qr_left = (prs.slide_width - qr_width) / 2
    qr_top = Inches(2.5)
    slide.shapes.add_picture(qr_bytes, qr_left, qr_top, width=qr_width)
    
    # Add registration code display
    code_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.8))
    code_frame = code_box.text_frame
    code_frame.word_wrap = True
    code_frame.text = f"Registration Code: {registration_code}"
    code_paragraph = code_frame.paragraphs[0]
    code_paragraph.font.size = Pt(28)
    code_paragraph.font.bold = True
    code_paragraph.font.color.rgb = RGBColor(30, 30, 30)
    code_paragraph.alignment = PP_ALIGN.CENTER
    code_paragraph.font.name = 'Courier New'
    
    # Add instructions
    instructions_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.6))
    instructions_frame = instructions_box.text_frame
    instructions_frame.word_wrap = True
    instructions_frame.text = "Point your phone camera at the QR code to register"
    instructions_paragraph = instructions_frame.paragraphs[0]
    instructions_paragraph.font.size = Pt(14)
    instructions_paragraph.font.color.rgb = RGBColor(100, 116, 139)
    instructions_paragraph.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_path = 'Lucky_Draw_QR_Code.pptx'
    prs.save(output_path)
    print(f"✅ PowerPoint file created: {output_path}")
    print(f"📱 QR Code contains: {registration_code}")
    return output_path

if __name__ == "__main__":
    # Generate with default code
    generate_qr_pptx("Z5G-12345")
    
    print("\n📝 To generate with a different code, use:")
    print("   from generate_qr_pptx import generate_qr_pptx")
    print("   generate_qr_pptx('ABC-54321')")
